#!/usr/bin/env python3
# -*- coding: utf-8 -*-

r"""render_report_pptx.py

Escape-safety:
- Regex patterns are raw strings (r"...")
- No invalid backslash escapes in docstrings; backslashes shown as "\\n".

Compatibility:
- Token-based templates
- Marker-based Templates Slides ("Titre section", "Texte Texte", "Légende Photo")

"""

from __future__ import annotations

import argparse
import json
import re
from copy import deepcopy
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Pt

from PIL import Image


def load_json(path: Path) -> Dict[str, Any]:
    return json.loads(path.read_text(encoding='utf-8'))


def normalize_whitespace(text: str) -> str:
    t = (text or '').replace('\r\n', '\n').replace('\r', '\n')
    t = re.sub(r"\n{3,}", "\n\n", t)
    return t.strip()


def strip_markdown(text: str) -> str:
    if not text:
        return ''
    t = text
    t = re.sub(r"\*\*(.*?)\*\*", r"\1", t)
    t = re.sub(r"\*(.*?)\*", r"\1", t)
    t = re.sub(r"^\s*#+\s*", "", t, flags=re.MULTILINE)
    t = t.replace('•', '-')
    return t.strip()


def sanitize_client_body(text: str) -> str:
    if not text:
        return ''
    t = normalize_whitespace(strip_markdown(text))
    out: List[str] = []
    for ln in t.splitlines():
        s = ln.strip()
        if not s:
            out.append('')
            continue
        low = s.lower()
        if 'page_id=' in low:
            continue
        if low.startswith('preuve:') or 'preuve:' in low:
            continue
        out.append(ln)
    return re.sub(r"\n{3,}", "\n\n", "\n".join(out)).strip()


def clone_slide(prs_out: Presentation, slide_in) -> Any:
    blank_layout = prs_out.slide_layouts[6]
    new_slide = prs_out.slides.add_slide(blank_layout)
    try:
        new_slide._element.get_or_add_bg()._set_bg(slide_in._element.get_or_add_bg())
    except Exception:
        pass
    for shape in slide_in.shapes:
        new_slide.shapes._spTree.insert_element_before(deepcopy(shape._element), 'p:extLst')
    return new_slide


def slide_text(slide) -> str:
    parts: List[str] = []
    for shape in slide.shapes:
        if getattr(shape, 'has_text_frame', False) and shape.has_text_frame:
            parts.append(shape.text_frame.text or '')
    return "\n".join(parts)


def replace_text_in_shape(shape, mapping: Dict[str, str]) -> None:
    if not getattr(shape, 'has_text_frame', False) or not shape.has_text_frame:
        return
    tf = shape.text_frame
    full = tf.text or ''
    replaced = full
    for k, v in mapping.items():
        if k in replaced:
            replaced = replaced.replace(k, v)
    if replaced != full:
        tf.text = replaced


def replace_placeholders(slide, mapping: Dict[str, str]) -> None:
    for shape in slide.shapes:
        replace_text_in_shape(shape, mapping)


def detect_template_slide_types(tpl: Presentation, slide_types_cfg: Dict[str, Any]) -> Dict[str, Any]:
    types = (slide_types_cfg.get('types') or {})
    found: Dict[str, Any] = {}
    for sl in tpl.slides:
        stxt = slide_text(sl)
        for tname, tcfg in types.items():
            toks = tcfg.get('detect_tokens') or []
            if toks and all(tok in stxt for tok in toks):
                found.setdefault(tname, sl)
    return found


def infer_repo_root(assembled_path: Path) -> Path:
    start = assembled_path.resolve()
    for cand in [start.parent] + list(start.parents):
        if (cand / 'process').exists() and (cand / 'input').exists():
            return cand
        if cand.name.lower() == 'process' and (cand.parent / 'input').exists():
            return cand.parent
    return assembled_path.parent.parent.parent


_img_cache: Dict[str, Optional[Path]] = {}
_conv_cache: Dict[Path, Path] = {}


def resolve_image_path(img: str, base_dir: Path, repo_root: Path) -> Optional[Path]:
    if not img:
        return None
    if img in _img_cache:
        return _img_cache[img]

    p = Path(img)
    if p.is_absolute() and p.exists():
        _img_cache[img] = p
        return p

    cand = (base_dir / img).resolve()
    if cand.exists():
        _img_cache[img] = cand
        return cand

    bn = p.name
    cand2 = (base_dir / bn).resolve()
    if cand2.exists():
        _img_cache[img] = cand2
        return cand2

    for r in (
        repo_root / 'process' / 'onenote',
        repo_root / 'input' / 'onenote-exporter',
        repo_root / 'input' / 'onenote-exporter' / 'output',
        repo_root / 'input',
    ):
        if not r.exists():
            continue
        c3 = r / img
        if c3.exists():
            _img_cache[img] = c3
            return c3
        c4 = r / bn
        if c4.exists():
            _img_cache[img] = c4
            return c4

    # slow fallback
    for sub in ('process', 'input'):
        base = repo_root / sub
        if not base.exists():
            continue
        try:
            for hit in base.rglob(bn):
                if hit.is_file():
                    _img_cache[img] = hit
                    return hit
        except Exception:
            continue

    _img_cache[img] = None
    return None


def magic_type(p: Path) -> str:
    try:
        b = p.read_bytes()[:16]
    except Exception:
        return 'unknown'
    if len(b) >= 3 and b[0:3] == b'\xFF\xD8\xFF':
        return 'jpeg'
    if len(b) >= 8 and b[0:8] == b'\x89PNG\r\n\x1a\n':
        return 'png'
    if len(b) >= 12 and b[0:4] == b'RIFF' and b[8:12] == b'WEBP':
        return 'webp'
    if len(b) >= 6 and (b[0:6] == b'GIF87a' or b[0:6] == b'GIF89a'):
        return 'gif'
    return 'unknown'


def normalize_image_for_ppt(p: Path, tmp_dir: Path) -> Optional[Path]:
    if not p.exists():
        return None
    if p in _conv_cache:
        return _conv_cache[p]
    kind = magic_type(p)
    ext = p.suffix.lower()
    if kind == 'jpeg' and ext in ('.jpg', '.jpeg'):
        _conv_cache[p] = p
        return p
    if kind == 'png' and ext == '.png':
        _conv_cache[p] = p
        return p
    try:
        im = Image.open(p)
        if im.mode not in ('RGB', 'RGBA'):
            im = im.convert('RGB')
        out = tmp_dir / (p.stem + '.png')
        im.save(out, format='PNG', optimize=True)
        _conv_cache[p] = out
        return out
    except Exception:
        return None


def images_to_pairs(images: List[Any]) -> List[Tuple[str, str]]:
    out: List[Tuple[str, str]] = []
    for im in images or []:
        if isinstance(im, str):
            out.append((im, ''))
        elif isinstance(im, dict):
            out.append((im.get('path') or '', im.get('caption') or ''))
    return [(p, c) for p, c in out if p]


def _shape_area(sh) -> int:
    try:
        return int(sh.width) * int(sh.height)
    except Exception:
        return 0


def _pic_has_image_rel(sh) -> bool:
    try:
        _ = sh.image
        return True
    except Exception:
        return False


def _is_picture_slot(sh) -> bool:
    try:
        if sh.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
            try:
                return 'PICTURE' in str(sh.placeholder_format.type).upper()
            except Exception:
                return False
    except Exception:
        pass
    try:
        return sh.shape_type == MSO_SHAPE_TYPE.PICTURE and not _pic_has_image_rel(sh)
    except Exception:
        return False


def pick_image_slots(slide, max_slots: int) -> List[Any]:
    pics = [sh for sh in slide.shapes if _is_picture_slot(sh)]
    # de-dup by rect
    seen = set(); uniq = []
    for sh in pics:
        key = (int(sh.left), int(sh.top), int(sh.width), int(sh.height))
        if key in seen:
            continue
        seen.add(key)
        uniq.append(sh)
    uniq.sort(key=_shape_area, reverse=True)
    return uniq[:max_slots]


def _crop_to_fill(pic_shape, img_path: Path, box_w: int, box_h: int) -> None:
    try:
        im = Image.open(img_path)
        w, h = im.size
        if not w or not h or not box_h:
            return
        img_ar = w / h
        box_ar = float(box_w) / float(box_h)
        pic_shape.crop_left = 0.0
        pic_shape.crop_right = 0.0
        pic_shape.crop_top = 0.0
        pic_shape.crop_bottom = 0.0
        if img_ar > box_ar:
            new_w = h * box_ar
            excess = max(0.0, w - new_w)
            frac = (excess / w) / 2.0 if w else 0.0
            pic_shape.crop_left = frac
            pic_shape.crop_right = frac
        else:
            new_h = w / box_ar
            excess = max(0.0, h - new_h)
            frac = (excess / h) / 2.0 if h else 0.0
            pic_shape.crop_top = frac
            pic_shape.crop_bottom = frac
    except Exception:
        return


def overlay_images(slide, images: List[Any], base_dir: Path, repo_root: Path, tmp_dir: Path, stats: Dict[str, Any]) -> None:
    pairs = images_to_pairs(images)
    if not pairs:
        return
    slots = pick_image_slots(slide, max_slots=len(pairs))
    for idx, (path, _) in enumerate(pairs[:len(slots)]):
        stats['requested'] += 1
        ip = resolve_image_path(path, base_dir, repo_root)
        if not ip:
            stats['unresolved'].append(path)
            continue
        normp = normalize_image_for_ppt(ip, tmp_dir)
        if not normp:
            stats['unresolved'].append(path)
            continue
        ph = slots[idx]
        pic = slide.shapes.add_picture(str(normp), ph.left, ph.top, width=ph.width, height=ph.height)
        _crop_to_fill(pic, normp, int(ph.width), int(ph.height))
        stats['embedded'] += 1


def fill_legends(slide, images: List[Any]) -> None:
    pairs = images_to_pairs(images)
    legend_shapes = []
    # exact then containing
    for shape in slide.shapes:
        if getattr(shape, 'has_text_frame', False) and shape.has_text_frame:
            txt = (shape.text_frame.text or '')
            if txt.strip() == 'Légende Photo' or 'Légende Photo' in txt:
                legend_shapes.append(shape)
    seen = set(); uniq = []
    for sh in legend_shapes:
        if id(sh) in seen:
            continue
        seen.add(id(sh)); uniq.append(sh)
    legend_shapes = uniq
    if not pairs:
        for sh in legend_shapes:
            try:
                sh.text_frame.text = ''
            except Exception:
                pass
        return
    for i, (_, cap) in enumerate(pairs[:len(legend_shapes)]):
        cap = (cap or '').strip()
        if len(cap) > 60:
            cap = cap[:57].rstrip() + '…'
        try:
            legend_shapes[i].text_frame.text = cap
        except Exception:
            pass


def set_bullets_fit(shape, body: str, *, max_lines: int = 12, min_font_pt: int = 10) -> None:
    if not getattr(shape, 'has_text_frame', False) or not shape.has_text_frame:
        return
    tf = shape.text_frame
    lines = [ln.strip() for ln in normalize_whitespace(strip_markdown(sanitize_client_body(body))).split('\n') if ln.strip()]
    out = []
    for ln in lines:
        s = ln
        while s.startswith('- '):
            s = s[2:].strip()
        if len(s) > 190:
            s = s[:187].rstrip() + '…'
        out.append(s)
    if len(out) > max_lines:
        out = out[:max_lines-1] + ['…']
    for p in tf.paragraphs:
        p.text = ''
    if not out:
        return
    for i, s in enumerate(out):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = s
        p.level = 0
        p.alignment = PP_ALIGN.LEFT
    tf.word_wrap = True
    try:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except Exception:
        pass
    try:
        for p in tf.paragraphs:
            if p.font.size is not None and p.font.size < Pt(min_font_pt):
                p.font.size = Pt(min_font_pt)
    except Exception:
        pass


def _set_title_any(slide_out, title: str) -> None:
    if not title:
        return
    for sh in slide_out.shapes:
        if getattr(sh, 'has_text_frame', False) and sh.has_text_frame:
            txt = sh.text_frame.text or ''
            if '{{SLIDE_TITLE}}' in txt or 'Titre section' in txt:
                sh.text_frame.text = title
                return


def _set_body_any(slide_out, body: str, *, max_lines: int) -> None:
    for sh in slide_out.shapes:
        if getattr(sh, 'has_text_frame', False) and sh.has_text_frame:
            txt = sh.text_frame.text or ''
            if '{{TEXTE_BULLETS}}' in txt or 'Texte Texte' in txt:
                set_bullets_fit(sh, body, max_lines=max_lines)
                return


def build_deck(base_template: Path, assembled_path: Path, out_path: Path, slide_types_path: Path, repo_root_override: Optional[Path] = None) -> None:
    base_tpl = Presentation(str(base_template))
    data = load_json(assembled_path)
    repo_root = repo_root_override or infer_repo_root(assembled_path)
    base_dir = assembled_path.parent
    tmp_dir = base_dir / '_tmp_img'
    tmp_dir.mkdir(parents=True, exist_ok=True)

    stats = {'requested': 0, 'embedded': 0, 'unresolved': []}
    cfg = load_json(slide_types_path) if slide_types_path.exists() else {'types': {}, 'defaults': {}}
    catalog = detect_template_slide_types(base_tpl, cfg)

    today = datetime.now().strftime('%d/%m/%Y')
    global_map = {
        '{{DATE}}': today,
        '{{VERSION}}': '1',
        '{{CONTACT_EMAIL}}': 'contact@build4use.eu',
    }

    prs_out = Presentation()
    prs_out.slide_width = base_tpl.slide_width
    prs_out.slide_height = base_tpl.slide_height

    def emit_content_slide(s: Dict[str, Any]):
        stype = (s.get('type') or 'CONTENT_TEXT').strip()
        title = (s.get('title') or '').strip()
        body = (s.get('bullets') or s.get('body') or '')
        images = s.get('images') or []

        slide_in = catalog.get(stype) or catalog.get('CONTENT_TEXT_IMAGES') or catalog.get('CONTENT_TEXT')
        if not slide_in:
            return
        slide_out = clone_slide(prs_out, slide_in)
        replace_placeholders(slide_out, global_map)
        _set_title_any(slide_out, title)
        _set_body_any(slide_out, body, max_lines=12 if images else 14)
        fill_legends(slide_out, images)
        overlay_images(slide_out, images, base_dir, repo_root, tmp_dir, stats)

    slides_spec = data.get('slides')
    if isinstance(slides_spec, list):
        for s in slides_spec:
            if isinstance(s, dict) and s.get('type') != 'PART_DIVIDER':
                emit_content_slide(s)

    prs_out.save(str(out_path))

    print(f"Repo root: {repo_root}")
    print(f"Images: embedded {stats['embedded']} / requested {stats['requested']}")
    if stats['unresolved']:
        print('Unresolved image paths (sample):')
        for u in stats['unresolved'][:12]:
            print(' -', u)


def main() -> None:
    ap = argparse.ArgumentParser()
    ap.add_argument('--template', required=True)
    ap.add_argument('--assembled', required=True)
    ap.add_argument('--out', default='Rapport_Audit.pptx')
    ap.add_argument('--slide-types', default='input/config/slide_types.json')
    ap.add_argument('--project-root', default='')
    args = ap.parse_args()
    root_override = Path(args.project_root) if args.project_root else None
    build_deck(Path(args.template), Path(args.assembled), Path(args.out), Path(args.slide_types), repo_root_override=root_override)


if __name__ == '__main__':
    main()
